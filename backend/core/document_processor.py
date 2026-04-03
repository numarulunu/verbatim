"""
Document Processor - Extract text from DOCX, XLSX, PPTX, TXT, RTF

Supports multiple document formats:
- DOCX (Microsoft Word)
- XLSX (Microsoft Excel)
- PPTX (Microsoft PowerPoint)
- TXT (Plain text)
- RTF (Rich Text Format)
- CSV (Comma-separated values)

Features:
- Table extraction from Word/Excel
- Slide extraction from PowerPoint
- Metadata preservation
- Error handling for corrupted files
"""

from pathlib import Path
from typing import Optional, Dict, Any, List
import os

from .quiet import quiet_print

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    import xlrd  # For legacy .xls files
    XLRD_AVAILABLE = True
except ImportError:
    XLRD_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from striprtf.striprtf import rtf_to_text
    RTF_AVAILABLE = True
except ImportError:
    RTF_AVAILABLE = False

import csv


class DocumentProcessingError(Exception):
    """Base exception for document processing issues"""


class MissingDependencyError(DocumentProcessingError):
    """Raised when a required optional dependency is missing"""


class UnsupportedFormatError(DocumentProcessingError):
    """Raised when an unsupported file format is requested"""


class DocumentProcessor:
    """
    Extract text from various document formats

    Supports:
    - DOCX: Paragraphs, tables, headers/footers
    - XLSX: All sheets, all cells
    - PPTX: All slides, titles, content
    - TXT: Plain text
    - RTF: Formatted text
    - CSV: All rows/columns
    """

    def __init__(self):
        """Initialize document processor"""
        self._check_dependencies()

    def _check_dependencies(self):
        """Check which format handlers are available"""
        if not DOCX_AVAILABLE:
            quiet_print("Warning: python-docx not available - DOCX processing disabled")
        if not XLSX_AVAILABLE:
            quiet_print("Warning: openpyxl not available - XLSX processing disabled")
        if not XLRD_AVAILABLE:
            quiet_print("Warning: xlrd<2.0 not available - legacy XLS processing disabled")
        if not PPTX_AVAILABLE:
            quiet_print("Warning: python-pptx not available - PPTX processing disabled")
        if not RTF_AVAILABLE:
            quiet_print("Warning: striprtf not available - RTF processing disabled")

    def extract_text(self, file_path: str) -> str:
        """
        Extract text from document

        Args:
            file_path: Path to document file

        Returns:
            Extracted text

        Raises:
            FileNotFoundError: If file doesn't exist
            ValueError: If format not supported
        """
        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        suffix = path.suffix.lower()

        # Route to appropriate handler
        if suffix == '.docx':
            return self._extract_docx(file_path)
        elif suffix == '.xlsx':
            return self._extract_xlsx(file_path)
        elif suffix == '.xls':
            return self._extract_xls(file_path)
        elif suffix == '.pptx':
            return self._extract_pptx(file_path)
        elif suffix == '.txt':
            return self._extract_txt(file_path)
        elif suffix == '.rtf':
            return self._extract_rtf(file_path)
        elif suffix == '.csv':
            return self._extract_csv(file_path)
        else:
            raise UnsupportedFormatError(f"Unsupported format: {suffix}")

    def _extract_docx(self, file_path: str) -> str:
        """
        Extract text from DOCX file

        Extracts:
        - All paragraphs
        - All tables (formatted as markdown)
        - Preserves basic structure
        """
        if not DOCX_AVAILABLE:
            raise MissingDependencyError("python-docx not installed, cannot process DOCX files")

        try:
            doc = Document(file_path)

            text_parts = []

            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # Extract tables
            for table_idx, table in enumerate(doc.tables):
                text_parts.append(f"\n[Table {table_idx + 1}]")

                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    text_parts.append(" | ".join(cells))

                text_parts.append("")  # Blank line after table

            return "\n".join(text_parts)

        except Exception as e:
            raise DocumentProcessingError(f"Error extracting DOCX: {e}") from e

    def _extract_xlsx(self, file_path: str) -> str:
        """
        Extract text from XLSX file

        Extracts:
        - All sheets
        - All cells with values
        - Preserves row structure
        """
        if not XLSX_AVAILABLE:
            raise MissingDependencyError("openpyxl not installed, cannot process XLSX files")

        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)

            text_parts = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"\n[Sheet: {sheet_name}]\n")

                for row in sheet.iter_rows(values_only=True):
                    # Filter out None values and convert to strings
                    cells = [str(cell) for cell in row if cell is not None]
                    if cells:  # Only add non-empty rows
                        text_parts.append(" | ".join(cells))

                text_parts.append("")  # Blank line between sheets

            return "\n".join(text_parts)

        except Exception as e:
            raise DocumentProcessingError(f"Error extracting XLSX: {e}") from e

    def _extract_xls(self, file_path: str) -> str:
        """
        Extract text from legacy XLS file using xlrd (<2.0)
        """
        if not XLRD_AVAILABLE:
            raise MissingDependencyError("xlrd<2.0 not installed, cannot process legacy XLS files")

        try:
            workbook = xlrd.open_workbook(file_path)
            text_parts: List[str] = []

            for sheet in workbook.sheets():
                text_parts.append(f"\n[Sheet: {sheet.name}]\n")
                for row_idx in range(sheet.nrows):
                    row_values = [
                        str(sheet.cell_value(row_idx, col)).strip()
                        for col in range(sheet.ncols)
                        if str(sheet.cell_value(row_idx, col)).strip()
                    ]
                    if row_values:
                        text_parts.append(" | ".join(row_values))

                text_parts.append("")

            return "\n".join(text_parts)

        except Exception as e:
            raise DocumentProcessingError(f"Error extracting XLS: {e}") from e

    def _extract_pptx(self, file_path: str) -> str:
        """
        Extract text from PPTX file

        Extracts:
        - All slides
        - Titles and content
        - Speaker notes
        """
        if not PPTX_AVAILABLE:
            raise MissingDependencyError("python-pptx not installed, cannot process PPTX files")

        try:
            prs = Presentation(file_path)

            text_parts = []

            for slide_idx, slide in enumerate(prs.slides):
                text_parts.append(f"\n[Slide {slide_idx + 1}]\n")

                # Extract all text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)

                # Extract notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes.strip():
                        text_parts.append(f"\n[Notes: {notes}]")

                text_parts.append("")  # Blank line between slides

            return "\n".join(text_parts)

        except Exception as e:
            raise DocumentProcessingError(f"Error extracting PPTX: {e}") from e

    def _extract_txt(self, file_path: str) -> str:
        """
        Extract text from TXT file

        Handles multiple encodings
        """
        encodings = ['utf-8', 'latin-1', 'cp1252', 'ascii']

        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue

        # Fallback: read as binary and decode with errors='ignore'
        try:
            with open(file_path, 'rb') as f:
                return f.read().decode('utf-8', errors='ignore')
        except Exception as e:
            raise DocumentProcessingError(f"Error extracting TXT: {e}") from e

    def _extract_rtf(self, file_path: str) -> str:
        """
        Extract text from RTF file
        """
        if not RTF_AVAILABLE:
            raise MissingDependencyError("striprtf not installed, cannot process RTF files")

        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()

            return rtf_to_text(rtf_content)

        except Exception as e:
            raise DocumentProcessingError(f"Error extracting RTF: {e}") from e

    def _extract_csv(self, file_path: str) -> str:
        """
        Extract text from CSV file
        """
        try:
            text_parts = []

            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)

                for row in reader:
                    # Filter empty cells
                    cells = [cell.strip() for cell in row if cell.strip()]
                    if cells:
                        text_parts.append(" | ".join(cells))

            return "\n".join(text_parts)

        except Exception as e:
            raise DocumentProcessingError(f"Error extracting CSV: {e}") from e

    def get_metadata(self, file_path: str) -> Dict[str, Any]:
        """
        Extract metadata from document

        Args:
            file_path: Path to document

        Returns:
            Dict with metadata (author, created, modified, etc.)
        """
        path = Path(file_path)

        metadata = {
            'filename': path.name,
            'extension': path.suffix.lower(),
            'size_bytes': path.stat().st_size,
            'modified': path.stat().st_mtime
        }

        # Format-specific metadata
        suffix = path.suffix.lower()

        try:
            if suffix == '.docx' and DOCX_AVAILABLE:
                doc = Document(file_path)
                core_props = doc.core_properties
                metadata.update({
                    'author': core_props.author or 'Unknown',
                    'created': core_props.created,
                    'modified': core_props.modified,
                    'title': core_props.title or '',
                    'pages': len(doc.paragraphs)
                })

            elif suffix == '.xlsx' and XLSX_AVAILABLE:
                wb = openpyxl.load_workbook(file_path)
                metadata.update({
                    'sheets': len(wb.sheetnames),
                    'sheet_names': wb.sheetnames
                })

            elif suffix == '.pptx' and PPTX_AVAILABLE:
                prs = Presentation(file_path)
                metadata.update({
                    'slides': len(prs.slides)
                })

        except Exception:
            pass  # Metadata extraction failed, continue with basic metadata

        return metadata

    def is_supported(self, file_path: str) -> bool:
        """
        Check if file format is supported

        Args:
            file_path: Path to file

        Returns:
            True if format is supported
        """
        suffix = Path(file_path).suffix.lower()
        supported = ['.docx', '.xlsx', '.xls', '.pptx', '.txt', '.rtf', '.csv']
        return suffix in supported


# Convenience function
def extract_document_text(file_path: str) -> str:
    """
    Extract text from document (convenience function)

    Args:
        file_path: Path to document

    Returns:
        Extracted text
    """
    processor = DocumentProcessor()
    return processor.extract_text(file_path)
